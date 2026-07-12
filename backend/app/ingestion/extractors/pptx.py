"""PPTX extraction via python-pptx: slide titles, body text, tables, and

speaker notes, positioned by slide number."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from app.ingestion.extractors.base import BaseExtractor, ExtractionError
from app.ingestion.models import ElementType, ExtractedDocument, ExtractedElement


class PptxExtractor(BaseExtractor):
    extensions = (".pptx",)
    file_type = "pptx"

    def extract(self, path: Path) -> ExtractedDocument:
        doc = self._new_document(path)
        try:
            prs = Presentation(str(path))
        except Exception as exc:
            raise ExtractionError(f"Cannot open PPTX {path.name}: {exc}") from exc

        props = prs.core_properties
        doc.native_properties = {
            "title": props.title or None,
            "author": props.author or None,
        }
        doc.page_count = len(prs.slides)

        for slide_no, slide in enumerate(prs.slides, start=1):
            title_shape = slide.shapes.title
            if title_shape is not None and title_shape.has_text_frame:
                title = title_shape.text_frame.text.strip()
                if title:
                    doc.elements.append(
                        ExtractedElement(ElementType.SLIDE_TITLE, title, position=slide_no)
                    )

            for shape in slide.shapes:
                if shape is title_shape:
                    continue
                if shape.has_table:
                    rows = [
                        " | ".join(cell.text.strip() for cell in row.cells)
                        for row in shape.table.rows
                    ]
                    doc.elements.append(
                        ExtractedElement(
                            ElementType.TABLE, "\n".join(rows), position=slide_no
                        )
                    )
                elif shape.has_text_frame:
                    text = "\n".join(
                        p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()
                    )
                    if text:
                        doc.elements.append(
                            ExtractedElement(ElementType.SLIDE_BODY, text, position=slide_no)
                        )

            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    doc.elements.append(
                        ExtractedElement(ElementType.SPEAKER_NOTES, notes, position=slide_no)
                    )
        return doc

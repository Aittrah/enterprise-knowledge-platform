"""Shared fixtures: build one real sample file per supported format."""

from __future__ import annotations

from pathlib import Path

import pytest


def minimal_pdf(text: str) -> bytes:
    """Hand-assemble a valid single-page PDF with a text layer."""
    stream = f"BT /F1 24 Tf 72 720 Td ({text}) Tj ET".encode()
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        (
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            b"/Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>"
        ),
        b"<< /Length %d >>\nstream\n%s\nendstream" % (len(stream), stream),
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    out = bytearray(b"%PDF-1.4\n")
    offsets = []
    for i, obj in enumerate(objects, 1):
        offsets.append(len(out))
        out += b"%d 0 obj\n" % i + obj + b"\nendobj\n"
    xref_pos = len(out)
    out += b"xref\n0 %d\n" % (len(objects) + 1)
    out += b"0000000000 65535 f \n"
    for off in offsets:
        out += b"%010d 00000 n \n" % off
    out += b"trailer\n<< /Size %d /Root 1 0 R >>\nstartxref\n%d\n%%%%EOF" % (
        len(objects) + 1,
        xref_pos,
    )
    return bytes(out)


@pytest.fixture
def sample_txt(tmp_path: Path) -> Path:
    path = tmp_path / "policy.txt"
    path.write_text(
        "Leave Policy\n\nEmployees accrue 22 days of annual leave.\n\n"
        "Unused leave carries over up to 5 days.",
        encoding="utf-8",
    )
    return path


@pytest.fixture
def sample_html(tmp_path: Path) -> Path:
    path = tmp_path / "handbook.html"
    path.write_text(
        """<html><head><title>Employee Handbook</title>
        <script>alert('skip me')</script></head><body>
        <h1>Benefits</h1>
        <p>Health insurance is provided from day one.</p>
        <ul><li>Dental cover</li><li>Vision cover</li></ul>
        <table><tr><th>Plan</th><th>Cost</th></tr>
        <tr><td>Basic</td><td>0</td></tr></table>
        </body></html>""",
        encoding="utf-8",
    )
    return path


@pytest.fixture
def sample_csv(tmp_path: Path) -> Path:
    path = tmp_path / "invoices.csv"
    path.write_text(
        "invoice_id,vendor,amount\nINV-001,Acme Corp,1200.50\nINV-002,Globex,89.99\n",
        encoding="utf-8",
    )
    return path


@pytest.fixture
def sample_docx(tmp_path: Path) -> Path:
    import docx

    document = docx.Document()
    document.core_properties.title = "Q3 Financial Report"
    document.core_properties.author = "Finance Team"
    document.add_heading("Executive Summary", level=1)
    document.add_paragraph("Revenue grew 14% quarter over quarter.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Revenue"
    table.cell(1, 1).text = "$2.4M"
    path = tmp_path / "report.docx"
    document.save(str(path))
    return path


@pytest.fixture
def sample_pptx(tmp_path: Path) -> Path:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # title + content
    slide.shapes.title.text = "Platform Overview"
    slide.placeholders[1].text = "Hybrid retrieval with citations"
    slide.notes_slide.notes_text_frame.text = "Mention GraphRAG here"
    path = tmp_path / "deck.pptx"
    prs.save(str(path))
    return path


@pytest.fixture
def sample_pdf(tmp_path: Path) -> Path:
    path = tmp_path / "contract.pdf"
    path.write_bytes(minimal_pdf("Hello EKIP PDF"))
    return path
